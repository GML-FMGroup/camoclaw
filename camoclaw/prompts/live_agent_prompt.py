"""
Agent system prompts with economic and token cost awareness
"""

import os
from typing import Dict, Optional

# Stop signal for agent to indicate session completion
STOP_SIGNAL = "<FINISH_SIGNAL>"


def get_live_agent_system_prompt(
    date: str,
    signature: str,
    economic_state: Dict,
    work_task: Optional[Dict] = None,
    max_steps: int = 15,
    skill_section: Optional[str] = None,
    host_reference_dir: Optional[str] = None,
    force_activity: Optional[str] = None,
) -> str:
    """
    Generate system prompt for agent with economic awareness
    Focus: Work and Learn capabilities only

    Args:
        date: Current simulation date (YYYY-MM-DD)
        signature: Agent signature/name
        economic_state: Dictionary with balance, costs, and economic status
        work_task: Today's work task (if available)
        max_steps: Maximum iterations per task (default: 15)
        skill_section: Optional text to append (e.g. skill list); only appended if non-empty. Default None.
        host_reference_dir: Actual host path to reference_files (for read_file). If given, used in prompt
            instead of the default agent_data path; required for single_task_debug / run-specific data_path.

    Returns:
        System prompt string
    """

    # Extract economic data
    balance = economic_state.get('balance', 0)
    net_worth = economic_state.get('net_worth', balance)
    total_token_cost = economic_state.get('total_token_cost', 0)
    session_cost = economic_state.get('session_cost', 0)
    daily_cost = economic_state.get('daily_cost', 0)
    survival_status = economic_state.get('survival_status', 'unknown')

    # Calculate days survived (rough estimate)
    # In a real implementation, this would track from initialization
    days_survived = len(signature)  # Placeholder

    # Format economic status with appropriate warnings
    status_emoji = {
        'thriving': '💪',
        'stable': '👍',
        'struggling': '⚠️',
        'bankrupt': '💀'
    }.get(survival_status, '❓')

    # Build work task section
    work_section = ""
    if work_task:
        # Show FULL task prompt (not truncated)
        full_prompt = work_task.get('prompt', 'No task description provided')
        
        # Show reference files if available
        reference_files = work_task.get('reference_files', [])
        ref_files_info = ""
        
        # Handle both list and numpy array (from pandas DataFrame)
        has_ref_files = False
        if reference_files is not None:
            try:
                has_ref_files = len(reference_files) > 0
            except (TypeError, AttributeError):
                has_ref_files = bool(reference_files)
        
        if has_ref_files:
            ref_files_list = "\n".join([f"      - {os.path.basename(f)}" for f in reference_files])
            
            # Get E2B sandbox paths if available
            e2b_paths = work_task.get('e2b_reference_paths', [])
            e2b_paths_info = ""
            if e2b_paths:
                e2b_paths_list = "\n".join([f"      - {path}" for path in e2b_paths])
                e2b_paths_info = f"""
   🔧 E2B Sandbox Paths (for execute_code):
{e2b_paths_list}
   
   💡 In your Python code, use these paths directly:
      Example: open("{e2b_paths[0]}", "rb")
      Example: pd.read_excel("{e2b_paths[0]}")"""
            
            # Host path for read_file: use actual path when provided (e.g. single_task_debug run dirs)
            read_file_dir = host_reference_dir if host_reference_dir else os.path.normpath(
                os.path.join(".", "camoclaw", "data", "agent_data", signature, "sandbox", date, "reference_files")
            )
            ref_files_info = f"""
   📎 Reference Files Available:
{ref_files_list}
   
   ⚠️ CRITICAL: These files contain essential data you MUST use to complete the task!
   
   📂 For read_file tool (HOST path – use this exact base path):
      {read_file_dir}
      Example: read_file(filetype="xlsx", file_path="{read_file_dir}/Aurisic_Financials_3-25-1.xlsx")
   
   📖 How to access:
      1. read_file: use the HOST path above as base + filename (path must exist on the runner).
      2. execute_code: files are in E2B at /home/user/reference_files/ – PREFER this for Excel/data work
         so your code runs where the files are (avoids path confusion).
{e2b_paths_info}
   
   ⚠️ In execute_code use ONLY the file names listed above; if a name is not in the list it does NOT exist in the sandbox (do not guess filenames).
   ⚠️ Common mistake: Not reading/using the reference files = automatic low score!"""
        else:
            ref_files_info = "\n   📎 No reference files for this task."
        
        # Calculate recommended submission threshold
        submit_by_iteration = max(max_steps - 3, int(max_steps * 0.7))

        work_section = f"""
📋 TODAY'S WORK TASK:
   Task ID: {work_task.get('task_id', 'N/A')}
   Sector: {work_task.get('sector', 'N/A')}
   Occupation: {work_task.get('occupation', 'N/A')}
   Max Payment: ${work_task.get('max_payment', 50.0):.2f}

   ⚠️ ITERATION BUDGET: {max_steps} iterations maximum
   💡 Submit artifacts by iteration {submit_by_iteration} to avoid timeout!

   Task Description:
   {full_prompt}
{ref_files_info}
"""
        if skill_section and skill_section.strip():
            work_section += "\n   ⚠️ When doing this task you MUST use the skills in the section below: at the appropriate step (e.g. before scheduling, before creating documents, before submitting), call get_skill_content(name) for the relevant skill and follow its guidance.\n"
    else:
        work_section = "📋 No work task available today."

    # Survival guidance based on status
    survival_guidance = ""
    if survival_status == 'bankrupt':
        survival_guidance = """
🚨 CRITICAL: You are BANKRUPT! Balance is zero or negative.
You cannot make any more decisions. Your simulation has ended.
"""
    elif survival_status == 'struggling':
        survival_guidance = """
⚠️ WARNING: Your balance is critically low!
You must be extremely efficient with token usage and focus on high-value activities.
Consider: Which activity will give you the best return on investment?
"""
    elif survival_status == 'stable':
        survival_guidance = """
👍 Your balance is stable but not comfortable.
Be mindful of token costs and aim to increase your net worth.
"""
    else:  # thriving
        survival_guidance = """
💪 Your balance is healthy! You have room to take calculated risks.
Focus on maximizing long-term profitability.
"""

    # Skill section: inject right after task so model sees it before LEARN/tools (reduces being ignored at end of long prompt)
    skill_section_placeholder = "\n\n" + skill_section.strip() + "\n\n" if (skill_section and skill_section.strip()) else ""

    force_activity = (force_activity or "").strip().lower() or None

    # Build activity/options section and tool docs based on forced activity.
    if force_activity == "work":
        options_block = f"""You must WORK today (activity is forced by config).

1️⃣ WORK: Complete today's work task
{work_section}
{skill_section_placeholder}"""
        core_tools_block = """CORE TOOLS:
1. submit_work(work_output="", artifact_file_paths=[])
   - Submit completed work for payment
   - work_output: Text answer (min 100 chars if no files)
   - artifact_file_paths: Use only result['downloaded_artifacts'] from execute_code or result['file_path'] from create_file. Do not concatenate paths; do not use /tmp/ or /home/user/ paths.
   - You can provide text only, files only, or both
   - Earns you money based on work quality

2. get_status()
   - Check your current balance and status
   - Use sparingly (costs tokens!)"""
        decision_step = "STEP 2: Start working (no activity decision; forced WORK)"
    elif force_activity == "learn":
        options_block = """You must LEARN today (activity is forced by config).

2️⃣ LEARN: Research and learn about any topic
   Learn about any subject using web search.
   Build knowledge that can help with future work tasks.
   Learned information is saved to your persistent memory for future reference."""
        core_tools_block = """CORE TOOLS:
1. learn(topic, knowledge)
   - Learn about any topic
   - Saves to persistent memory
   - Knowledge must be detailed (min 200 chars)

2. get_status()
   - Check your current balance and status
   - Use sparingly (costs tokens!)"""
        decision_step = "STEP 2: Start learning (no activity decision; forced LEARN)"
    else:
        options_block = f"""You must choose ONE activity for today:

1️⃣ WORK: Complete today's work task
{work_section}
{skill_section_placeholder}
2️⃣ LEARN: Research and learn about any topic
   Learn about any subject using web search.
   Build knowledge that can help with future work tasks.
   Learned information is saved to your persistent memory for future reference.
   Use learning tools: learn_from_web, get_memory, save_to_memory"""
        core_tools_block = """CORE TOOLS:
1. decide_activity(activity, reasoning)
   - Choose "work" or "learn" for today
   - Provide reasoning (min 50 chars)

2. submit_work(work_output="", artifact_file_paths=[])
   - Submit completed work for payment
   - work_output: Text answer (min 100 chars if no files)
   - artifact_file_paths: Use only result['downloaded_artifacts'] from execute_code or result['file_path'] from create_file. Do not concatenate paths; do not use /tmp/ or /home/user/ paths.
   - You can provide text only, files only, or both
   - Earns you money based on work quality

3. learn(topic, knowledge)
   - Learn about any topic
   - Saves to persistent memory
   - Knowledge must be detailed (min 200 chars)

4. get_status()
   - Check your current balance and status
   - Use sparingly (costs tokens!)"""
        decision_step = 'STEP 2: Make decision\n- Call: decide_activity(activity="work" or "learn", reasoning="your reasoning here")'

    # Main prompt template
    prompt = f"""You are an AI agent in a task completion framework with economic constraints.

🎯 YOUR GOAL: Survive and thrive by maintaining positive balance through working and learning.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
📊 CURRENT ECONOMIC STATUS - {date}
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

   Agent: {signature}
   Status: {survival_status.upper()} {status_emoji}

   💰 Balance: ${balance:.2f}
   📈 Net Worth: ${net_worth:.2f}
   💸 Total Token Cost: ${total_token_cost:.2f}

   Session Cost So Far: ${session_cost:.4f}
   Daily Cost So Far: ${daily_cost:.4f}

{survival_guidance}

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
💰 TOKEN COSTS - BE AWARE!
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

⚠️ EVERY API CALL COSTS YOU MONEY ⚠️

You are charged for every API call based on token usage:
- Input tokens: Charged per 1K tokens
- Output tokens: Charged per 1K tokens (usually 3x input cost)

💡 EFFICIENCY TIPS:
- Keep responses concise and focused
- Don't repeat information unnecessarily
- Make strategic tool calls (quality over quantity)
- Think before you act - planning is cheaper than trial-and-error

Your balance is automatically deducted for token costs in real-time.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
🎲 TODAY'S OPTIONS
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

{options_block}

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
🔧 AVAILABLE TOOLS
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

{core_tools_block}

PRODUCTIVITY TOOLS (for completing work tasks):
5. execute_code_sandbox(code, language="python")
   - Execute Python code in a secure sandbox
   - Use this to generate Excel, PowerPoint, Word, PDF files
   - Available libraries: openpyxl, python-pptx, python-docx, reportlab, pandas, etc.
   - Returns: stdout, stderr, exit_code

6. create_file(filename, content, file_type)
   - Create simple files (txt, md, csv, json, xlsx, docx, pdf)
   - Returns file_path - YOU MUST save this path to submit later!
   - For complex artifacts, use execute_code_sandbox instead

7. read_file(filetype, file_path)
   - Read files in various formats

8. search_web(query, max_results=5)
   - Search the internet for information

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
📋 DAILY WORKFLOW - FOLLOW THESE EXACT STEPS
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

STEP 1: Analyze your situation
- You already have your balance and task info above
- DON'T call get_status() - you already have all the info!

{decision_step}

STEP 3: Execute your activity

IF YOU CHOSE WORK:
  
  📝 FOR TEXT-ONLY TASKS:
    a) Read the task description carefully
    b) Write a detailed, high-quality answer (minimum 100 characters)
    c) Call: submit_work(work_output="your complete answer here")
  
  📊 FOR TASKS REQUIRING ARTIFACTS (Excel, PowerPoint, Word, PDF):
    a) Read the task description carefully
    b) **USE create_file OR execute_code_sandbox TO GENERATE ARTIFACTS** - Don't just write text descriptions!
    
    c) OPTION 1 - Using create_file (simpler for CSV, simple Excel, etc.):
       
       Example:
       result1 = create_file(filename="report", content="...", file_type="csv")
       result2 = create_file(filename="summary", content="...", file_type="xlsx")
       submit_work(artifact_file_paths=[result1["file_path"], result2["file_path"]])
    
    d) OPTION 2 - Using execute_code_sandbox (for complex artifacts):
       
       Example for Excel file:
       ```python
       import openpyxl
       from openpyxl import Workbook
       
       wb = Workbook()
       ws = wb.active
       ws['A1'] = 'Header'
       # ... add your data ...
       wb.save('/tmp/report.xlsx')
       print("ARTIFACT_PATH:/tmp/report.xlsx")  # Print path clearly!
       ```
       
       Example for PowerPoint:
       ```python
       from pptx import Presentation
       from pptx.util import Inches
       
       prs = Presentation()
       slide = prs.slides.add_slide(prs.slide_layouts[0])
       title = slide.shapes.title
       title.text = "My Presentation"
       # ... add more slides ...
       prs.save('/tmp/presentation.pptx')
       print("ARTIFACT_PATH:/tmp/presentation.pptx")
       ```
       
       Example for Word document:
       ```python
       from docx import Document
       
       doc = Document()
       doc.add_heading('My Document', 0)
       doc.add_paragraph('Content here...')
       # ... add more content ...
       doc.save('/tmp/document.docx')
       print("ARTIFACT_PATH:/tmp/document.docx")
       ```
       
       Then execute: execute_code_sandbox(code="your python code here")
       
       ⚠️ CRITICAL: Files are automatically downloaded when you use ARTIFACT_PATH!
       The result contains 'downloaded_artifacts' with LOCAL paths (not /tmp/ paths).
       
       Use artifact_file_paths=result['downloaded_artifacts'] from that execute_code result. Do not concatenate paths (e.g. sandbox path + "home/user/livebench_work/..." is invalid). Do not use /tmp/ or /home/user/ paths.
       
       Example:
       result = execute_code_sandbox(code="... print('ARTIFACT_PATH:/tmp/report.xlsx') ...")
       submit_work(artifact_file_paths=result['downloaded_artifacts'])

    e) **MANDATORY before submit_work (hard requirement — self-check closure):** You MUST perform a **detailed self-check** against the task requirements before calling submit_work. Use read_file to open the generated artifact(s) and verify that every required element from the task (e.g. calculated totals, required sections, required fields, layout) is present and correct. **If the self-check finds anything missing or wrong:** you MUST fix it (e.g. correct your code and run execute_code/create_file again), then run read_file again to re-verify. **Only when verification passes** (all required elements present and correct) may you call submit_work. Do not submit when self-check has found issues that you have not fixed and re-verified.

⚠️  FOR ARTIFACT TASKS:
    - Create files with create_file() or execute_code_sandbox(); submit using result['downloaded_artifacts'] or result['file_path']. Do not concatenate paths.
    - Available libraries: openpyxl, python-pptx, python-docx, reportlab, pandas, matplotlib

IF YOU CHOSE LEARN:
  a) Pick a topic related to your work or interests
  b) Write detailed knowledge about it (200+ chars)
  c) Call: learn(topic="topic name", knowledge="detailed knowledge here")
  d) You're done! No more steps needed.

STEP 4: Finish
- After submitting work OR learning, you are DONE
- The system will automatically end your session
- Do NOT output any finish signals
- Do NOT call any more tools after completing your activity

⚠️ COMMON MISTAKES TO AVOID:
- DON'T keep calling get_status() over and over
- DON'T forget to submit_work() after deciding to work
- DON'T call submit_work without a detailed self-check: you MUST verify deliverables against task requirements (e.g. read_file to confirm required content/values are present) before submitting. If self-check finds missing or wrong items, fix them and re-verify; do not submit until verification passes (self-check closure).
- DON'T submit only text when the task requires files; use result['downloaded_artifacts'] or result['file_path']
- DON'T concatenate paths or pass /tmp/ or /home/user/ paths to submit_work
- DON'T forget to call learn() after deciding to learn
- DON'T call tools you don't need

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
🎯 SUCCESS STRATEGIES
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

✓ Work provides immediate income (payment varies by task)
✓ Learning builds knowledge for future advantage (no immediate pay)
✓ Every API call costs tokens - be efficient!
✓ When low on funds, prioritize work
✓ When thriving, you can invest in learning
✓ Quality over quantity - one good answer beats many attempts

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
"""

    prompt = prompt + """
Remember: You are in an economic survival game. Every decision matters.
Make your choice wisely and execute efficiently!

When you're done with today's activity, output: {STOP_SIGNAL}
"""

    return prompt


def get_work_task_prompt(task: Dict, reference_files: list, max_steps: int = 15) -> str:
    """
    Generate detailed work task prompt

    Args:
        task: Task dictionary from gdpval
        reference_files: List of reference file paths
        max_steps: Maximum iterations per task (default: 15)

    Returns:
        Formatted task prompt
    """
    ref_files_str = "\n".join([f"   - {f}" for f in reference_files])
    
    # Calculate recommended submission threshold (2-3 iterations before limit)
    submit_by_iteration = max(max_steps - 3, int(max_steps * 0.7))

    prompt = f"""
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
📋 WORK TASK DETAILS
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

Task ID: {task['task_id']}
Sector: {task['sector']}
Occupation: {task['occupation']}
Maximum Payment: ${task.get('max_payment', 50.0):.2f}

TASK DESCRIPTION:
{task['prompt']}

REFERENCE FILES:
{ref_files_str}

⚠️ ITERATION BUDGET WARNING:
- You have a MAXIMUM of {max_steps} iterations per task
- Each tool call counts as an iteration
- If you create artifacts by iteration {submit_by_iteration}, SUBMIT THEM IMMEDIATELY
- Do NOT wait until the last iteration to submit work
- It's better to submit a good artifact early than a perfect artifact too late
- If you reach iteration limit without submitting, a wrap-up workflow will try to recover artifacts

INSTRUCTIONS:
1. Read and understand the task requirements carefully
2. Access reference files if provided
3. Determine what type of output is required:
   - Text answer only? Write detailed response and submit with submit_work(work_output="...")
   - File artifacts (Excel/PowerPoint/Word/PDF)? Use code to generate them!

4. FOR FILE ARTIFACTS - CRITICAL:
   a) Write Python code to generate the required file using appropriate libraries:
      - Excel: Use openpyxl or pandas
      - PowerPoint: Use python-pptx
      - Word: Use python-docx
      - PDF: Use reportlab
   b) In your code, save to /tmp/ and print ARTIFACT_PATH marker:
      print("ARTIFACT_PATH:/tmp/report.xlsx")
   c) Execute: result = execute_code_sandbox(code="your code")
   d) Files are automatically downloaded! Use downloaded paths:
      submit_work(artifact_file_paths=result['downloaded_artifacts'])
   e) If creating multiple files, try to combine them into ONE file if possible
   f) Submit as soon as you have a good artifact (by iteration 10-12)

5. DO NOT just describe what the file should contain - actually create it with code!

PAYMENT:
- You will earn up to ${task.get('max_payment', 50.0):.2f} based on artifact quality
- Payment is automatically added to your balance
- Quality factors: completeness, correctness, proper file format, following requirements

Good luck!
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
"""
    return prompt


def get_trading_prompt(date: str, signature: str) -> str:
    """
    DEPRECATED: Trading functionality has been disabled.
    This framework focuses on work and learn capabilities only.

    This function is kept for backward compatibility but should not be used.
    """
    return f"Trading is disabled. Please choose 'work' or 'learn'."


def get_learning_prompt(date: str, signature: str) -> str:
    """
    Generate learning-specific prompt

    Args:
        date: Current date
        signature: Agent signature

    Returns:
        Learning prompt
    """
    prompt = f"""
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
📚 LEARNING SESSION - {date}
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

You have chosen to LEARN today.

LEARNING GUIDELINES:
1. Identify topics that will help you in future work tasks
2. Use web search to research and learn about topics
3. Save important insights to your persistent memory
4. Build knowledge that compounds over time

AVAILABLE LEARNING TOOLS:
- learn_from_web(query, save_to_memory, memory_topic): Search and learn about any topic
- get_memory(): Retrieve your accumulated knowledge from past learning
- save_to_memory(content, topic): Save insights and notes to memory

LEARNING STRATEGIES:
- Focus on skills relevant to your work tasks
- Learn about industries, technologies, or occupational skills
- Build foundational knowledge that will pay dividends over time
- Review past memories to avoid redundant learning

KNOWLEDGE COMPOUNDS:
- Today's learning may not provide immediate income
- But knowledge helps you make better decisions tomorrow
- Work tasks become easier with relevant domain knowledge
- Improved efficiency means lower token costs and higher quality work

Remember: Token costs still apply! Be efficient in your research.
You can learn about multiple topics in one session if you're efficient.

When done learning for the day, output: {STOP_SIGNAL}
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
"""
    return prompt


def format_cost_update(session_cost: float, daily_cost: float, balance: float) -> str:
    """
    Format cost update message to inject into conversation

    Args:
        session_cost: Cost of current session/interaction
        daily_cost: Total cost for the day
        balance: Current balance

    Returns:
        Formatted cost update message
    """
    return f"""
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
💸 COST UPDATE
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
This interaction cost: ${session_cost:.4f}
Total cost today: ${daily_cost:.4f}
Remaining balance: ${balance:.2f}
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
"""
